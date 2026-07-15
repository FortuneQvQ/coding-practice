from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "校园资讯聚合系统_答辩PPT模板.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

GREEN_DARK = RGBColor(18, 78, 55)
GREEN = RGBColor(28, 120, 78)
GREEN_MID = RGBColor(38, 154, 99)
GREEN_LIGHT = RGBColor(216, 243, 220)
GREEN_PALE = RGBColor(241, 251, 244)
TEXT = RGBColor(39, 55, 48)
MUTED = RGBColor(105, 120, 112)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(116, 178, 135)
SHADOW = RGBColor(207, 226, 214)
ACCENT = RGBColor(246, 177, 66)
ACCENT_DARK = RGBColor(194, 116, 32)


def set_text(run, size=20, color=TEXT, bold=False):
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_textbox(slide, x, y, w, h, text, size=22, color=TEXT, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    set_text(run, size, color, bold)
    return box


def add_round_rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.1)
    return shape


def add_pill(slide, x, y, w, text, fill=GREEN, color=WHITE, size=12, height=0.38):
    pill = add_round_rect(slide, x, y, w, height, fill, fill)
    tf = pill.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_text(run, size, color, True)
    return pill


def add_title(slide, title, subtitle=None, section="Campus News Hub"):
    add_pill(slide, 0.72, 0.38, 1.65, section, GREEN_DARK)
    marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(0.92), Inches(0.08), Inches(0.42))
    marker.fill.solid()
    marker.fill.fore_color.rgb = ACCENT
    marker.line.fill.background()
    add_textbox(slide, 0.98, 0.88, 8.14, 0.48, title, 24, GREEN_DARK, True)
    if subtitle:
        add_textbox(slide, 0.72, 1.31, 11.9, 0.24, subtitle, 11, MUTED, False, PP_ALIGN.CENTER)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(1.76), Inches(11.9), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN
    line.line.fill.background()


def add_footer(slide, idx):
    add_textbox(slide, 0.72, 7.03, 5.8, 0.2, "校园资讯聚合系统", 8.5, MUTED)
    add_textbox(slide, 12.25, 7.0, 0.45, 0.25, f"{idx:02d}", 9, GREEN_DARK, True, PP_ALIGN.RIGHT)


def blank_slide(prs, idx, title=None, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = GREEN_PALE
    top_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.12))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = GREEN_DARK
    top_bar.line.fill.background()
    if title:
        add_title(slide, title, subtitle)
    add_footer(slide, idx)
    return slide


def add_bullets(slide, x, y, w, h, items, size=15, color=TEXT, bullet=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(6)
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        if bullet:
            p._p.get_or_add_pPr().set("marL", "285750")
            p._p.get_or_add_pPr().set("indent", "-171450")
    return box


def add_card(slide, x, y, w, h, title, body, tag=None, compact=False):
    shadow = add_round_rect(slide, x + 0.05, y + 0.06, w, h, SHADOW, SHADOW)
    card = add_round_rect(slide, x, y, w, h, WHITE, LINE)
    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.08))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = GREEN
    stripe.line.fill.background()
    if tag:
        add_pill(slide, x + 0.22, y + 0.2, 0.78, tag, GREEN_DARK, WHITE, 10.5, 0.33)
        title_y = y + 0.68
    else:
        title_y = y + 0.28
    title_x = x + 0.26
    title_w = w - 0.52
    title_size = 13
    if compact:
        title_x = x + 0.24
        title_y = y + 0.66
        title_w = 1.28
        title_size = 11.5
    add_textbox(slide, title_x, title_y, title_w, 0.32, title, title_size, GREEN_DARK, True)
    add_bullets(slide, x + 0.28, title_y + 0.48, w - 0.56, h - 0.7, body, 11.5, MUTED, False)
    return card


def add_placeholder(slide, x, y, w, h, title, hint):
    box = add_round_rect(slide, x, y, w, h, RGBColor(250, 253, 251), LINE)
    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.08))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = GREEN_DARK
    stripe.line.fill.background()
    add_textbox(slide, x + 0.2, y + h / 2 - 0.2, w - 0.4, 0.28, title, 14, GREEN_DARK, True, PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.35, y + h / 2 + 0.17, w - 0.7, 0.28, hint, 10, MUTED, False, PP_ALIGN.CENTER)
    return box


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = GREEN_PALE

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(2.65))
    band.fill.solid()
    band.fill.fore_color.rgb = GREEN_DARK
    band.line.fill.background()

    add_pill(slide, 0.82, 0.62, 1.9, "Campus News Hub", ACCENT, GREEN_DARK)
    add_textbox(slide, 0.8, 1.55, 9.6, 0.8, "校园资讯聚合系统", 42, WHITE, True)
    add_textbox(slide, 0.86, 2.42, 8.8, 0.38, "校园信息一站式获取与检索平台", 18, ACCENT, True)
    add_textbox(slide, 0.9, 3.05, 7.2, 0.36, "是啊吃什么组 · 小组项目答辩", 15, MUTED)

    add_round_rect(slide, 8.7, 2.1, 3.45, 2.45, WHITE, LINE)
    add_textbox(slide, 9.04, 2.5, 2.8, 0.35, "一站式校园信息平台", 18, GREEN_DARK, True, PP_ALIGN.CENTER)
    add_bullets(slide, 9.15, 3.05, 2.65, 0.95, ["新闻聚合", "分类浏览", "关键词搜索"], 13, MUTED)

    for i, (label, desc) in enumerate([
        ("Crawler", "资讯采集"),
        ("Parser", "字段提取"),
        ("SQLite", "本地检索"),
        ("Website", "页面展示"),
    ]):
        x = 1.0 + i * 2.75
        add_card(slide, x, 5.1, 2.2, 0.9, label, [], None)
        add_textbox(slide, x + 0.26, 5.72, 1.68, 0.2, desc, 10.5, MUTED, False, PP_ALIGN.CENTER)
        if i < 3:
            add_textbox(slide, x + 2.28, 5.37, 0.3, 0.2, "→", 18, ACCENT_DARK)

    add_textbox(slide, 0.9, 6.78, 5.5, 0.25, "汇报人：XXX / XXX / XXX    指导教师：XXX", 10.5, MUTED)
    add_textbox(slide, 11.85, 6.78, 0.5, 0.25, "01", 9, GREEN_DARK, True, PP_ALIGN.RIGHT)
    return slide


def agenda(prs, idx):
    slide = blank_slide(prs, idx, "汇报结构", "从需求出发，展示实现路径，并回到系统价值与可扩展方向")
    items = [
        ("总", "项目定位与建设目标", "为什么要做：减少多站点浏览成本，避免遗漏关键信息"),
        ("分", "系统设计与模块实现", "怎么做：爬取、解析、入库、生成网站与交互检索"),
        ("总", "项目成果与后续规划", "做成什么：跑通核心链路，并预留持续演进空间"),
    ]
    for i, (tag, title, desc) in enumerate(items):
        y = 2.35 + i * 1.22
        add_pill(slide, 1.05, y + 0.1, 0.55, tag, GREEN if i != 1 else GREEN_DARK)
        add_textbox(slide, 1.85, y, 3.8, 0.35, title, 20, GREEN_DARK, True)
        add_textbox(slide, 1.86, y + 0.47, 8.6, 0.25, desc, 13, MUTED)
    add_placeholder(slide, 9.4, 2.4, 2.65, 2.85, "核心目标", "让校园资讯从多站点分散浏览，变为一个入口集中获取")


def project_overview(prs, idx):
    slide = blank_slide(prs, idx, "项目概述", "面向校园通知、新闻与活动的信息聚合系统，服务学生高效获取资讯")
    add_card(slide, 0.9, 2.1, 3.6, 2.0, "真实痛点", ["多站点浏览成本高", "重要通知容易遗漏"], "WHY")
    add_card(slide, 4.85, 2.1, 3.6, 2.0, "建设目标", ["采集并统一结构化资讯", "提供浏览、检索与筛选"], "GOAL")
    add_card(slide, 8.8, 2.1, 3.6, 2.0, "项目边界", ["单机 C++17 + SQLite", "不含登录与移动端功能"], "SCOPE")
    add_placeholder(slide, 1.15, 4.8, 10.8, 1.05, "项目定位", "将分散的校园资讯自动汇聚为可浏览、可检索、可追溯的统一入口")


def functions_slide(prs, idx):
    slide = blank_slide(prs, idx, "功能总览", "覆盖“采集—处理—存储—展示”闭环，兼顾学生使用与维护扩展")
    labels = [
        ("多来源采集", "按配置抓取公开资讯"),
        ("结构化解析", "提取标题、时间等字段"),
        ("关键词检索", "定位相关资讯"),
        ("浏览与筛选", "按来源和时间查阅"),
        ("本地持久化", "SQLite 全文检索"),
        ("静态页面生成", "首页、详情与搜索页"),
    ]
    for i, (title, body) in enumerate(labels):
        x = 0.95 + (i % 3) * 3.95
        y = 2.05 + (i // 3) * 1.75
        add_card(slide, x, y, 3.3, 1.18, title, [], f"0{i+1}", compact=True)
        add_textbox(slide, x + 1.56, y + 0.54, 1.48, 0.26, body, 9.5, MUTED, False, PP_ALIGN.CENTER)
    add_placeholder(slide, 2.0, 5.95, 9.3, 0.58, "首页浏览 → 来源/关键词检索 → 详情查看原文", "")


def architecture_slide(prs, idx):
    slide = blank_slide(prs, idx, "系统架构", "从数据采集到页面展示的完整处理流水线")
    steps = [
        ("配置源", "config.json\n站点 URL 列表"),
        ("爬虫模块", "libcurl 抓取\nmanifest 记录"),
        ("解析模块", "Tidy + pugixml\nJSON 字段提取"),
        ("数据库", "SQLite\nNews 结构体"),
        ("网站生成", "模板替换\nHTML/JSON 输出"),
        ("用户访问", "首页/分类\n搜索/详情"),
    ]
    for i, (title, body) in enumerate(steps):
        x = 0.7 + i * 2.08
        add_round_rect(slide, x, 2.45, 1.62, 1.35, WHITE, LINE)
        add_textbox(slide, x + 0.1, 2.72, 1.42, 0.25, title, 13, GREEN_DARK, True, PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.12, 3.08, 1.38, 0.4, body, 9.5, MUTED, False, PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_textbox(slide, x + 1.7, 2.95, 0.35, 0.28, "→", 17, GREEN)
    add_card(slide, 1.0, 4.85, 5.2, 1.22, "答辩表达", ["这页用来说明我们的项目不是孤立页面，而是一个可重复运行的数据生产与展示流程。"], None)
    add_card(slide, 6.75, 4.85, 5.2, 1.22, "可替换材料", ["可在流程图下方加入代码截图、数据库表截图或运行日志截图。"], None)


def module_slides(prs, idx):
    titles = [
        ("数据采集模块", "Crawler/base_01.cpp", ["读取 config.json 中的多站点 URL", "使用 libcurl 设置 UA、超时与自动跳转", "抓取结果保存为 page_x.html，并写入 manifest", "加入日志记录与 1 秒限流，便于调试和复现"]),
        ("解析与清洗模块", "Parser/HTMLparsing.cpp", ["使用 Tidy 将 HTML 清理为 XHTML", "使用 pugixml 与 XPath 提取链接和正文段落", "从页面 JSON 中解析标题、发布时间、正文内容", "将相对链接转换为绝对链接，减少后续访问错误"]),
        ("数据库模块", "DataBase/database.cpp", ["使用 SQLite 保存新闻结构化字段", "News 结构体统一承载标题、时间、来源、摘要等信息", "提供按时间、来源、主题、标题检索接口", "为网站生成与搜索分类功能提供数据支撑"]),
        ("网站生成模块", "Server/WebsiteGenerator.cpp", ["读取 HTML 模板并替换占位符", "生成首页、详情页、搜索页、分类首页与分类结果页", "输出 news.json 供前端脚本读取", "整体风格采用浅绿色渐变、白色卡片和圆角按钮"]),
    ]
    for offset, (title, file, bullets) in enumerate(titles):
        slide = blank_slide(prs, idx + offset, title, file)
        add_bullets(slide, 0.95, 2.15, 5.55, 3.45, bullets, 15, TEXT)
        add_placeholder(slide, 7.1, 2.0, 4.65, 3.15, "代码 / 运行截图占位", f"建议放：{file}")
        add_card(slide, 7.1, 5.55, 4.65, 0.85, "讲解重点", ["从输入、处理、输出三个角度讲清楚模块责任。"], None)


def frontend_slide(prs, idx):
    slide = blank_slide(prs, idx, "界面风格与交互", "延续网站浅绿色视觉，突出校园资讯的清爽与可读性")
    add_card(slide, 0.9, 2.1, 3.2, 1.45, "视觉语言", ["背景：#d8f3dc → #ffffff 渐变", "主色：#2d6a4f / #40916c", "内容：白色卡片 + 轻阴影"], "UI")
    add_card(slide, 4.45, 2.1, 3.2, 1.45, "首页交互", ["搜索框直接输入关键词", "最新新闻卡片展示摘要", "详情按钮进入新闻正文"], "HOME")
    add_card(slide, 8.0, 2.1, 3.2, 1.45, "分类交互", ["按时间、来源、主题分组", "标签按钮跳转结果页", "返回按钮形成闭环浏览"], "FLOW")
    add_placeholder(slide, 1.0, 4.35, 4.9, 1.65, "首页截图", "放置 Server/output/index.html")
    add_placeholder(slide, 6.45, 4.35, 4.9, 1.65, "分类页截图", "放置 Server/output/category_index.html")


def data_slide(prs, idx):
    slide = blank_slide(prs, idx, "数据结构与检索", "围绕 News 结构体完成从入库到页面展示的数据统一")
    fields = ["id", "title", "time", "content", "abstract", "url", "source", "image", "topic"]
    for i, field in enumerate(fields):
        x = 0.85 + (i % 5) * 2.35
        y = 2.1 + (i // 5) * 0.75
        add_pill(slide, x, y, 1.55, field, GREEN_LIGHT, GREEN_DARK)
    add_card(slide, 0.9, 4.05, 3.5, 1.35, "时间检索", ["今日新闻、最近一周、最近一个月、历史资讯"], None)
    add_card(slide, 4.85, 4.05, 3.5, 1.35, "来源检索", ["教务处、学生会、学院通知、学生社团等"], None)
    add_card(slide, 8.8, 4.05, 3.5, 1.35, "主题检索", ["教学通知、校园活动、竞赛比赛、考试安排、科研创新"], None)


def highlights_slide(prs, idx):
    slide = blank_slide(prs, idx, "技术亮点", "突出项目实现中最能体现工作量和工程性的部分")
    add_card(slide, 0.95, 2.0, 3.4, 2.55, "自动化链路", ["从 URL 配置到页面生成形成闭环", "减少手工维护新闻内容的成本", "模块之间通过结构化数据解耦"], "01")
    add_card(slide, 4.85, 2.0, 3.4, 2.55, "解析鲁棒性", ["Tidy 清理非标准 HTML", "pugixml/XPath 处理正文与链接", "相对路径转绝对路径"], "02")
    add_card(slide, 8.75, 2.0, 3.4, 2.55, "展示可扩展", ["模板化生成多类页面", "JSON 输出支持前端搜索与分类", "浅绿色校园风格统一"], "03")
    add_placeholder(slide, 2.25, 5.25, 8.6, 0.7, "答辩提示", "每个亮点都对应一段代码或一个页面效果，讲解时避免只说概念")


def demo_slide(prs, idx):
    slide = blank_slide(prs, idx, "演示路线", "建议按用户视角完成 2 分钟演示")
    route = [
        ("1", "打开首页", "展示标题、搜索框、最新资讯卡片"),
        ("2", "查看详情", "进入新闻详情页，说明字段完整性"),
        ("3", "分类浏览", "选择时间 / 来源 / 主题分类"),
        ("4", "关键词搜索", "输入关键词并查看结果"),
    ]
    for i, (num, title, desc) in enumerate(route):
        y = 2.0 + i * 0.95
        add_pill(slide, 0.95, y, 0.5, num, GREEN_DARK)
        add_textbox(slide, 1.72, y - 0.02, 2.3, 0.3, title, 17, GREEN_DARK, True)
        add_textbox(slide, 4.1, y + 0.02, 6.9, 0.25, desc, 13, MUTED)
    add_placeholder(slide, 7.35, 5.55, 4.2, 0.8, "演示链接 / 文件路径", "Server/output/index.html")


def summary_slide(prs, idx):
    slide = blank_slide(prs, idx, "总结与展望", "项目跑通校园资讯从采集到展示的核心链路，并具备持续扩展基础")
    add_card(slide, 0.95, 2.0, 3.45, 2.2, "核心成果", ["跑通采集、解析、存储与展示", "支持浏览、搜索与筛选"], "NOW")
    add_card(slide, 4.85, 2.0, 3.45, 2.2, "项目价值", ["建立统一校园资讯入口", "降低信息遗漏风险"], "VALUE")
    add_card(slide, 8.75, 2.0, 3.45, 2.2, "后续规划", ["定时更新与变更检测", "推送、推荐与桌面端适配"], "NEXT")
    add_placeholder(slide, 2.05, 5.25, 8.9, 0.7, "让重要通知更及时可达", "")


def thanks_slide(prs, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = GREEN_DARK
    add_textbox(slide, 0.9, 2.35, 11.4, 0.85, "谢谢聆听", 46, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, 0.9, 3.35, 11.4, 0.36, "欢迎老师和同学提问", 18, ACCENT, True, PP_ALIGN.CENTER)
    add_textbox(slide, 0.9, 4.15, 11.4, 0.3, "校园资讯聚合系统 · Campus News Hub", 13, GREEN_LIGHT, False, PP_ALIGN.CENTER)
    add_footer(slide, idx)


def rebalance_body_content(prs):
    desired_shift = Inches(0.18)
    body_top = Inches(1.9)
    footer_safe_top = Inches(6.72)
    for slide_index, slide in enumerate(prs.slides, start=1):
        if slide_index in (1, 15):
            continue
        body_shapes = [
            shape for shape in slide.shapes
            if body_top <= shape.top < Inches(6.8)
        ]
        if not body_shapes:
            continue
        lowest_edge = max(shape.top + shape.height for shape in body_shapes)
        shift = max(0, min(desired_shift, footer_safe_top - lowest_edge))
        for shape in body_shapes:
            shape.top += shift


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    cover(prs)
    agenda(prs, 2)
    project_overview(prs, 3)
    functions_slide(prs, 4)
    architecture_slide(prs, 5)
    module_slides(prs, 6)
    frontend_slide(prs, 10)
    data_slide(prs, 11)
    highlights_slide(prs, 12)
    demo_slide(prs, 13)
    summary_slide(prs, 14)
    thanks_slide(prs, 15)
    rebalance_body_content(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
